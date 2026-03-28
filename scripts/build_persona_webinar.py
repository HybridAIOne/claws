#!/usr/bin/env python3
"""Render HybridClaw persona slides into a PowerPoint deck.

This script reads the persona handbook DOCX, extracts the case sections,
renders a slide image for each persona plus a cover and closing slide, and
builds a PPTX with those images as full-slide artwork.
"""

from __future__ import annotations

import argparse
import io
import os
import re
import textwrap
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Inches


DOCX_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}


def repo_root() -> Path:
    return Path(__file__).resolve().parents[1]


def load_docx_paragraphs(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as zf:
        root = ET.fromstring(zf.read("word/document.xml"))
    paras: list[str] = []
    for p in root.findall(".//w:p", DOCX_NS):
        parts = [t.text for t in p.findall(".//w:t", DOCX_NS) if t.text]
        if parts:
            text = "".join(parts).strip()
            if text:
                paras.append(text)
    return paras


@dataclass
class PersonaCase:
    index: int
    code: str
    name: str
    role: str
    summary: str
    soul: list[str]
    identity: list[str]
    skills: list[str]
    cv: list[str]
    handover: str


def split_case_sections(paragraphs: list[str]) -> tuple[list[PersonaCase], list[str]]:
    case_starts = [i for i, line in enumerate(paragraphs) if line.startswith("Case ")]
    update_start = next(
        (i for i, line in enumerate(paragraphs) if line.startswith("Aktualisierung ")),
        len(paragraphs),
    )

    cases: list[PersonaCase] = []
    for pos, start in enumerate(case_starts):
        end = case_starts[pos + 1] if pos + 1 < len(case_starts) else update_start
        block = paragraphs[start:end]
        if len(block) < 10:
            continue

        m = re.match(r"Case\s+(\d+)\s+–\s+(.+)", block[0])
        index = int(m.group(1)) if m else len(cases) + 1
        name_line = block[1]
        name_match = re.match(r"\d+\.\s+(.+)", name_line)
        name = name_match.group(1) if name_match else name_line
        code = re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")
        role = block[2]
        summary = block[3]

        def section_lines(anchor: str, next_anchor: str | None = None) -> list[str]:
            try:
                a = block.index(anchor)
            except ValueError:
                return []
            if next_anchor and next_anchor in block[a + 1 :]:
                b = block.index(next_anchor, a + 1)
            else:
                b = len(block)
            return block[a + 1 : b]

        soul = section_lines("SOUL.md", "IDENTITY.md")
        identity = section_lines("IDENTITY.md", "manifest.json + workspace/skills/")
        skills = section_lines("manifest.json + workspace/skills/", "CV.md")
        cv = section_lines("CV.md", None)
        handover = next((line for line in block if line.startswith("Empfohlener Handover an Menschen:")), "")

        cases.append(
            PersonaCase(
                index=index,
                code=code,
                name=name,
                role=role,
                summary=summary,
                soul=soul,
                identity=identity,
                skills=skills,
                cv=cv,
                handover=handover.replace("Empfohlener Handover an Menschen:", "").strip(),
            )
        )

    update_lines = paragraphs[update_start:]
    return cases, update_lines


def find_font(candidates: Iterable[str], default: str | None = None) -> str:
    for candidate in candidates:
        if os.path.exists(candidate):
            return candidate
    if default:
        return default
    raise FileNotFoundError("No usable font file found")


FONT_REGULAR = find_font(
    [
        "/System/Library/Fonts/Avenir Next.ttc",
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/System/Library/Fonts/Arial.ttf",
    ]
)
FONT_MONO = find_font(
    [
        "/System/Library/Fonts/Menlo.ttc",
        "/System/Library/Fonts/Monaco.ttf",
    ],
    default=FONT_REGULAR,
)


def font(size: int, bold: bool = False, mono: bool = False):
    if mono:
        return ImageFont.truetype(FONT_MONO, size=size)
    index = 8 if bold else 7
    if FONT_REGULAR.endswith(".ttc"):
        return ImageFont.truetype(FONT_REGULAR, size=size, index=index)
    return ImageFont.truetype(FONT_REGULAR, size=size)


def measure(draw: ImageDraw.ImageDraw, text: str, fnt) -> tuple[int, int]:
    box = draw.textbbox((0, 0), text, font=fnt, stroke_width=0)
    return box[2] - box[0], box[3] - box[1]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt, width: int) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines() or [""]:
        if not raw_line.strip():
            lines.append("")
            continue
        words = raw_line.split()
        current = words[0]
        for word in words[1:]:
            trial = current + " " + word
            if measure(draw, trial, fnt)[0] <= width:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    fnt,
    fill,
    width: int,
    line_gap: int = 6,
) -> int:
    x, y = xy
    total_h = 0
    for line in wrap_text(draw, text, fnt, width):
        if line:
            draw.text((x, y), line, font=fnt, fill=fill)
            h = measure(draw, line, fnt)[1]
            y += h + line_gap
            total_h += h + line_gap
        else:
            y += line_gap + 4
            total_h += line_gap + 4
    return total_h


def draw_bullets(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    bullets: Iterable[str],
    width: int,
    color,
    bullet_color,
    title_font,
    body_font,
    gap: int = 10,
) -> int:
    cur = y
    for bullet in bullets:
        bullet = re.sub(r"^\s*[•\-\u2022]+\s*", "", bullet).strip()
        draw.ellipse((x, cur + 10, x + 8, cur + 18), fill=bullet_color)
        lines = wrap_text(draw, bullet, body_font, width - 20)
        first = True
        for line in lines:
            draw.text((x + 18, cur), line, font=body_font, fill=color)
            cur += measure(draw, line, body_font)[1] + (4 if first else 2)
            first = False
        cur += gap
    return cur - y


def paper_background(size: tuple[int, int], texture_path: Path | None = None) -> Image.Image:
    width, height = size
    base = Image.new("RGBA", size, "#f8f5ef")
    draw = ImageDraw.Draw(base)

    # Subtle paper grain.
    noise = Image.new("L", size, 0)
    px = noise.load()
    for y in range(height):
        for x in range(width):
            px[x, y] = 244 + ((x * 17 + y * 13) % 9)
    noise = noise.filter(ImageFilter.GaussianBlur(radius=0.8))
    grain = Image.new("RGBA", size, (255, 255, 255, 0))
    grain.putalpha(noise.point(lambda v: 16 if v else 0))
    base = Image.alpha_composite(base, grain)

    if texture_path and texture_path.exists():
        texture = Image.open(texture_path).convert("RGBA").resize(size, Image.Resampling.LANCZOS)
        texture = texture.convert("L").convert("RGBA")
        texture = texture.filter(ImageFilter.GaussianBlur(radius=18))
        alpha = Image.new("L", size, 0)
        alpha_px = alpha.load()
        for y in range(height):
            for x in range(width):
                alpha_px[x, y] = 16 + ((x * 7 + y * 11) % 6)
        texture.putalpha(alpha.point(lambda v: min(22, v)))
        base = Image.alpha_composite(base, texture)

    draw = ImageDraw.Draw(base)

    grid_color = (187, 190, 196, 38)
    for x in range(40, width, 160):
        draw.line((x, 32, x, height - 32), fill=grid_color, width=1)
    for y in range(30, height, 128):
        draw.line((32, y, width - 32, y), fill=grid_color, width=1)

    border_color = (170, 174, 182, 115)
    draw.rectangle((24, 24, width - 24, height - 24), outline=border_color, width=1)
    tick_color = (140, 148, 160, 180)
    for x in [24, width // 2, width - 24]:
        draw.line((x, 10, x, 40), fill=tick_color, width=1)
        draw.line((x, height - 40, x, height - 10), fill=tick_color, width=1)
    for y in [24, height // 2, height - 24]:
        draw.line((10, y, 40, y), fill=tick_color, width=1)
        draw.line((width - 40, y, width - 10, y), fill=tick_color, width=1)
    return base


def rounded_box(
    img: Image.Image,
    rect: tuple[int, int, int, int],
    fill,
    outline,
    radius: int = 22,
    shadow: bool = True,
    width: int = 3,
):
    x1, y1, x2, y2 = rect
    if shadow:
        shadow_img = Image.new("RGBA", img.size, (0, 0, 0, 0))
        shadow_draw = ImageDraw.Draw(shadow_img)
        shadow_draw.rounded_rectangle((x1 + 6, y1 + 8, x2 + 6, y2 + 8), radius=radius, fill=(33, 43, 59, 28))
        shadow_img = shadow_img.filter(ImageFilter.GaussianBlur(radius=6))
        img.alpha_composite(shadow_img)
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle(rect, radius=radius, fill=fill, outline=outline, width=width)


def draw_logo(img: Image.Image, logo_path: Path | None):
    if logo_path and logo_path.exists():
        im = Image.open(logo_path).convert("RGBA")
        max_w = 130
        ratio = max_w / im.width
        im = im.resize((max_w, max(1, int(im.height * ratio))), Image.Resampling.LANCZOS)
        img.alpha_composite(im, (img.width - im.width - 34, img.height - im.height - 24))
        return

    # Fallback wordmark when the template logo is unavailable.
    draw = ImageDraw.Draw(img)
    draw.text(
        (img.width - 252, img.height - 54),
        "hybridclaw.io",
        font=font(18, bold=True),
        fill=(14, 33, 62, 255),
    )


def title_block(img: Image.Image, title: str, subtitle: str | None = None, tag: str | None = None):
    draw = ImageDraw.Draw(img)
    tfont = font(58, bold=True)
    sfont = font(22, bold=False)
    tagfont = font(18, bold=True)
    x = 72
    y = 58
    if tag:
        pill_w = measure(draw, tag, tagfont)[0] + 36
        draw.rounded_rectangle((x, y - 10, x + pill_w, y + 28), radius=18, fill=(229, 223, 212, 255), outline=(167, 132, 94, 255), width=2)
        draw.text((x + 18, y - 3), tag, font=tagfont, fill=(103, 74, 43, 255))
        y += 58
    draw.text((x, y), title, font=tfont, fill=(15, 33, 62, 255))
    y += measure(draw, title, tfont)[1] + 12
    if subtitle:
        draw_wrapped(draw, (x, y), subtitle, sfont, (85, 95, 110, 255), img.width - 240, line_gap=5)
    draw.line((x, 202, img.width - 82, 202), fill=(180, 184, 190, 180), width=2)


def section_card(
    img: Image.Image,
    rect: tuple[int, int, int, int],
    header: str,
    header_color,
    body_lines: list[str],
    body_title: str | None = None,
    emphasize: bool = False,
):
    rounded_box(img, rect, fill=(253, 252, 249, 255), outline=(253, 252, 249, 255), radius=24, shadow=True)
    draw = ImageDraw.Draw(img)
    x1, y1, x2, y2 = rect
    header_h = 58
    draw.rounded_rectangle((x1, y1, x2, y1 + header_h + 1), radius=24, fill=header_color)
    draw.rectangle((x1, y1 + header_h - 1, x2, y1 + header_h + 2), fill=header_color)
    draw.rounded_rectangle((x1, y1, x2, y2), radius=24, outline=(66, 81, 106, 230), width=3)
    hfont = font(22, bold=True)
    draw.text((x1 + 20, y1 + 16), header, font=hfont, fill=(14, 33, 62, 255))

    body_x = x1 + 20
    body_y = y1 + header_h + 18
    body_font = font(19, bold=False)
    body_color = (44, 55, 71, 255)
    if body_title:
        title_font = font(20, bold=True)
        title_h = draw_wrapped(
            draw,
            (body_x, body_y),
            body_title,
            title_font,
            (15, 33, 62, 255),
            x2 - body_x - 20,
            line_gap=4,
        )
        body_y += title_h + 10
    y = body_y
    bullet_color = (166, 77, 64, 255)
    if body_lines:
        y += draw_bullets(draw, body_x, y, body_lines, x2 - body_x - 20, body_color, bullet_color, None, body_font, gap=8)


def render_cover(
    personas: list[PersonaCase],
    update_lines: list[str],
    out_path: Path,
    logo_path: Path | None,
    texture_path: Path | None,
):
    img = paper_background((1600, 900), texture_path=texture_path)
    draw = ImageDraw.Draw(img)
    title_block(
        img,
        "HybridClaw Persona Teambook",
        "AI coworker personas translated from the current claw directories and handbook.",
        tag="Updated 2026-03-25",
    )

    cfont = font(28, bold=True)
    sfont = font(18, bold=False)
    # Central constellation.
    center = (1040, 505)
    draw.rounded_rectangle((930, 370, 1150, 640), radius=30, fill=(20, 42, 79, 255), outline=(13, 24, 45, 255), width=4)
    draw.multiline_text((983, 445), "HybridClaw\nPersona Pack", font=font(32, bold=True), fill=(246, 244, 239, 255), align="center", spacing=6)
    draw.text((983, 560), f"{len(personas)} personas + 2 additions", font=font(20, bold=False), fill=(205, 211, 223, 255))

    positions = [
        (585, 285), (760, 250), (1230, 275), (1360, 430),
        (1240, 710), (910, 760), (620, 705), (430, 470),
        (500, 590), (760, 780), (1360, 610), (1080, 220),
    ]
    accents = [
        (192, 153, 90, 255), (166, 77, 64, 255), (110, 124, 142, 255), (192, 153, 90, 255),
        (166, 77, 64, 255), (110, 124, 142, 255), (192, 153, 90, 255), (166, 77, 64, 255),
        (110, 124, 142, 255), (192, 153, 90, 255), (166, 77, 64, 255), (110, 124, 142, 255),
    ]
    for idx, (persona, pos, accent) in enumerate(zip(personas, positions, accents), 1):
        x, y = pos
        draw.line((center[0], center[1], x + 110, y + 45), fill=(120, 131, 148, 170), width=3)
        draw.rounded_rectangle((x, y, x + 220, y + 76), radius=18, fill=(250, 248, 243, 255), outline=accent, width=3)
        draw.text((x + 16, y + 13), f"{persona.index:02d}", font=font(20, bold=True), fill=accent)
        draw.text((x + 58, y + 10), persona.name, font=font(18, bold=True), fill=(14, 33, 62, 255))
        draw.text((x + 58, y + 36), persona.role, font=font(15, bold=False), fill=(85, 95, 110, 255))

    # Text strip on the left.
    left_x = 72
    draw.rounded_rectangle((70, 250, 710, 690), radius=28, fill=(248, 246, 241, 240), outline=(183, 187, 194, 200), width=3)
    draw.text((98, 282), "What the pack standardizes", font=font(24, bold=True), fill=(15, 33, 62, 255))
    bullets = [
        "SOUL.md defines how each persona thinks, speaks, and decides.",
        "IDENTITY.md frames the work context, triggers, and handover points.",
        "manifest.json points to bundled, imported, and external skills.",
        "CV.md turns the role into a believable backstory and signature move.",
    ]
    draw_bullets(draw, 98, 336, bullets, 560, (44, 55, 71, 255), (166, 77, 64, 255), None, font(20), gap=14)

    draw.rounded_rectangle((70, 712, 710, 820), radius=22, fill=(17, 31, 54, 255), outline=(17, 31, 54, 255))
    draw.text((98, 740), "Outputs stay aligned with the current claws, not a static persona sheet.", font=font(22, bold=True), fill=(247, 244, 237, 255))
    draw.text((98, 776), "Current spec: manifest.skills + workspace/skills + vetted external git imports.", font=font(17), fill=(194, 200, 210, 255))

    draw_logo(img, logo_path)
    img.save(out_path)


def render_persona_slide(persona: PersonaCase, out_path: Path, logo_path: Path | None, texture_path: Path | None):
    img = paper_background((1600, 900), texture_path=texture_path)
    draw = ImageDraw.Draw(img)
    title_block(img, f"{persona.index:02d}. {persona.name}", None, tag=persona.role)

    # hero chips
    chip_y = 230
    chips = [
        ("SOUL", (195, 163, 94, 255)),
        ("IDENTITY", (110, 124, 142, 255)),
        ("SKILLS", (166, 77, 64, 255)),
        ("CV", (76, 90, 118, 255)),
    ]
    x = 74
    for label, color in chips:
        w = measure(draw, label, font(16, bold=True))[0] + 30
        draw.rounded_rectangle((x, chip_y, x + w, chip_y + 34), radius=17, fill=(250, 247, 240, 255), outline=color, width=2)
        draw.text((x + 15, chip_y + 8), label, font=font(16, bold=True), fill=color)
        x += w + 12

    # left and center cards
    section_card(
        img,
        (74, 282, 500, 720),
        "SOUL.md",
        (232, 224, 210, 255),
        persona.soul,
        body_title="Mission, tone, principles, guardrail",
    )
    section_card(
        img,
        (530, 282, 964, 720),
        "IDENTITY.md",
        (232, 224, 210, 255),
        persona.identity + [f"Recommended handover: {persona.handover}"] if persona.handover else persona.identity,
        body_title="Context, operating mode, and who they serve",
    )
    section_card(
        img,
        (1000, 282, 1526, 720),
        "manifest.json + workspace/skills/",
        (232, 224, 210, 255),
        persona.skills,
        body_title="Built-ins, ClawHub fit, connector fit",
    )

    # CV strip
    draw.rounded_rectangle((74, 744, 1526, 828), radius=22, fill=(18, 32, 56, 255), outline=(18, 32, 56, 255))
    draw.text((98, 765), "CV snapshot", font=font(18, bold=True), fill=(198, 205, 215, 255))
    cv_text = "  |  ".join(line.replace("• ", "") for line in persona.cv)
    draw_wrapped(draw, (218, 760), cv_text, font(18), (247, 244, 237, 255), 1260, line_gap=2)

    draw_logo(img, logo_path)
    img.save(out_path)


def render_closing_slide(
    update_lines: list[str],
    out_path: Path,
    logo_path: Path | None,
    texture_path: Path | None,
):
    img = paper_background((1600, 900), texture_path=texture_path)
    draw = ImageDraw.Draw(img)
    title_block(
        img,
        "Updated for the current claw directories.",
        "The teambook now reflects manifest-based skills, bundled skill directories, and the newest persona additions.",
        tag="Appendix / Update",
    )

    # Two featured additions.
    card_specs = [
        (
            (82, 282, 732, 700),
            "Leonie Graf",
            "LLM Evaluation Architect",
            [
                "Uses OpenAI-compatible APIs or browser chat surfaces.",
                "Turns topic knowledge into targeted LLM evaluation tasks.",
                "Useful when you need a repeatable benchmark, not a vibe check.",
            ],
            (110, 124, 142, 255),
        ),
        (
            (822, 282, 1520, 700),
            "Ronan Kade",
            "Codebase Specification Architect",
            [
                "Converts vague bug reports and feature requests into actionable specs.",
                "Grounds output in the actual repository and implementation constraints.",
                "Useful for tightening handoffs before engineering starts coding.",
            ],
            (166, 77, 64, 255),
        ),
    ]
    for rect, name, role, bullets, accent in card_specs:
        rounded_box(img, rect, fill=(252, 250, 245, 255), outline=accent, radius=26, width=3)
        x1, y1, x2, y2 = rect
        draw.rounded_rectangle((x1, y1, x2, y1 + 70), radius=26, fill=(232, 236, 241, 255), outline=accent, width=3)
        draw.text((x1 + 26, y1 + 20), name, font=font(28, bold=True), fill=(15, 33, 62, 255))
        draw.text((x1 + 26, y1 + 48), role, font=font(18, bold=False), fill=(86, 95, 108, 255))
        draw_bullets(draw, x1 + 26, y1 + 112, bullets, x2 - x1 - 52, (43, 54, 71, 255), accent, None, font(20), gap=12)

    draw.rounded_rectangle((82, 740, 1520, 822), radius=20, fill=(17, 31, 54, 255), outline=(17, 31, 54, 255))
    draw.text((108, 764), "Current rule of thumb", font=font(18, bold=True), fill=(198, 205, 215, 255))
    update_text = [
        "Skills now live in manifest.json as bundled, imports, or external git refs.",
        "Bundled skills belong under workspace/skills/ and are validated during packing.",
        "Use the teambook as a source of truth when building or exporting a claw.",
    ]
    draw_wrapped(draw, (300, 760), "  ".join(update_text), font(18), (247, 244, 237, 255), 1180, line_gap=2)

    draw_logo(img, logo_path)
    img.save(out_path)


def build_pptx(images: list[Path], out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for image_path in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    # Remove the default first slide if present.
    if len(prs.slides) > len(images):
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[0])
    prs.save(out_path)


def extract_template_logo(template_path: Path, cache_dir: Path) -> Path | None:
    if not template_path.exists():
        return None
    cache_dir.mkdir(parents=True, exist_ok=True)
    candidate = cache_dir / "template-logo.png"
    if candidate.exists():
        return candidate
    try:
        with zipfile.ZipFile(template_path) as zf:
            data = zf.read("ppt/media/image3.png")
    except Exception:
        return None
    candidate.write_bytes(data)
    return candidate


def extract_template_texture(template_path: Path, cache_dir: Path) -> Path | None:
    if not template_path.exists():
        return None
    cache_dir.mkdir(parents=True, exist_ok=True)
    candidate = cache_dir / "template-texture.png"
    if candidate.exists():
        return candidate
    for media_name in ("ppt/media/image4.png", "ppt/media/image5.png", "ppt/media/image6.png"):
        try:
            with zipfile.ZipFile(template_path) as zf:
                data = zf.read(media_name)
        except Exception:
            continue
        candidate.write_bytes(data)
        return candidate
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--docx",
        default="/Users/bkoehler/Downloads/HybridClaw_Teambook_AI_Coworker_Personas.docx",
        help="Path to the persona handbook docx.",
    )
    parser.add_argument(
        "--template",
        default=str(Path("/tmp/HybridClaw_Enterprise_AI_Assistant_Webinar.pptx")),
        help="Reference webinar pptx (used only for logo extraction conventions).",
    )
    parser.add_argument(
        "--out",
        default=str(Path("/tmp/HybridClaw_Enterprise_AI_Assistant_Webinar.generated.pptx")),
        help="Output pptx path.",
    )
    parser.add_argument(
        "--slides-dir",
        default=str(Path("/tmp/hybridclaw-persona-slides")),
        help="Directory where rendered slide PNGs will be written.",
    )
    args = parser.parse_args()

    docx_path = Path(args.docx)
    template_path = Path(args.template)
    out_path = Path(args.out)
    slides_dir = Path(args.slides_dir)
    slides_dir.mkdir(parents=True, exist_ok=True)

    paras = load_docx_paragraphs(docx_path)
    cases, update_lines = split_case_sections(paras)
    if len(cases) != 12:
        raise SystemExit(f"Expected 12 cases from docx, found {len(cases)}")

    logo_path = extract_template_logo(template_path, slides_dir)
    texture_path = extract_template_texture(template_path, slides_dir)
    slide_paths: list[Path] = []
    cover_path = slides_dir / "slide01_cover.png"
    render_cover(cases, update_lines, cover_path, logo_path, texture_path)
    slide_paths.append(cover_path)
    for idx, case in enumerate(cases, 2):
        p = slides_dir / f"slide{idx:02d}_{case.code}.png"
        render_persona_slide(case, p, logo_path, texture_path)
        slide_paths.append(p)
    closing_path = slides_dir / "slide14_closing.png"
    render_closing_slide(update_lines, closing_path, logo_path, texture_path)
    slide_paths.append(closing_path)

    build_pptx(slide_paths, out_path)
    print(out_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
